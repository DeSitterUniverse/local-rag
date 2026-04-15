import os
import sys
import time
import uuid
import json
import hashlib
import sqlite3
import uvicorn
import lancedb
import docx
import csv
import pptx
import openpyxl
import pyarrow as pa
import onnxruntime as ort
import numpy as np
import gc
from transformers import AutoTokenizer
from llama_cpp import Llama

def load_architecture_context() -> str:
    try:
        if getattr(sys, 'frozen', False):
            # Bundled environment (PyInstaller extracts to _MEIPASS)
            base_dir = sys._MEIPASS
        else:
            # Dev environment (One level up from python/)
            base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        
        target = os.path.join(base_dir, "AI_SYSTEM_AWARENESS.md")
        with open(target, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        return f"[Error loading internal architecture specs: {e}]"

ARCHITECTURE_CONTEXT = load_architecture_context()

from fastapi import FastAPI, BackgroundTasks, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from contextlib import asynccontextmanager
from pydantic import BaseModel
from pypdf import PdfReader
from langchain_text_splitters import RecursiveCharacterTextSplitter

DB_PATH = os.path.expanduser("~/cephalon-data")
os.makedirs(DB_PATH, exist_ok=True)
MODEL_DIR = os.path.expanduser("~/cephalon-data/models")
os.makedirs(MODEL_DIR, exist_ok=True)

# Define PyArrow Schema explicitly for LanceDB vectors
schema = pa.schema([
    pa.field("vector", pa.list_(pa.float32(), 768)),
    pa.field("id", pa.string()),
    pa.field("doc_id", pa.string()),
    pa.field("text", pa.string())
])

os.environ['HF_HOME'] = os.path.expanduser("~/.cephalon/models")
def _init_db(conn):
    """
    Initializes the SQLite metadata tracking database.
    While LanceDB handles the actual machine learning embeddings, SQLite acts as the DAG source of truth
    for UI state, tracking file upload status, deletion hooks, and chunk counts to prevent orphan vectors.
    """
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS documents (
            id TEXT PRIMARY KEY,
            path TEXT NOT NULL,
            content_hash TEXT NOT NULL,
            ingested_at INTEGER,
            chunk_count INTEGER,
            status TEXT DEFAULT 'pending',
            type TEXT DEFAULT 'file'
        );
        CREATE TABLE IF NOT EXISTS chunks (
            id TEXT PRIMARY KEY,
            doc_id TEXT REFERENCES documents(id) ON DELETE CASCADE,
            chunk_index INTEGER,
            text TEXT NOT NULL
        );
    """)
    try:
        conn.execute("ALTER TABLE documents ADD COLUMN type TEXT DEFAULT 'file'")
    except sqlite3.OperationalError:
        pass 
        
    cursor = conn.cursor()
    cursor.execute("INSERT OR IGNORE INTO documents (id, path, content_hash, chunk_count, status, type) VALUES (?, ?, ?, ?, ?, ?)", 
                   ("core_memory", "Internal AI Memory", "none", 0, "ready", "memory"))
    conn.commit()

@asynccontextmanager
async def lifespan(app: FastAPI):
    # Pure ONNX Boot — auto-extract bundled models on first run
    onnx_path = os.path.expanduser("~/cephalon-data/models/cross-encoder")
    embed_path = os.path.expanduser("~/cephalon-data/models/embedder")
    model_file = os.path.join(onnx_path, "model.onnx")
    embed_file = os.path.join(embed_path, "model.onnx")
    
    # If models don't exist locally, extract from bundled PyInstaller package
    if not os.path.exists(model_file) or not os.path.exists(embed_file):
        if getattr(sys, 'frozen', False):
            bundled_base = os.path.join(sys._MEIPASS, "onnx_models")
            bundled_cross = os.path.join(bundled_base, "cross-encoder")
            bundled_embed = os.path.join(bundled_base, "embedder")
            
            if os.path.exists(bundled_cross) and os.path.exists(bundled_embed):
                import shutil
                print("First boot detected — extracting bundled ONNX models...")
                if not os.path.exists(onnx_path):
                    shutil.copytree(bundled_cross, onnx_path)
                if not os.path.exists(embed_path):
                    shutil.copytree(bundled_embed, embed_path)
                print("ONNX models extracted successfully.")
            else:
                print("CRITICAL ERROR: Bundled ONNX models not found in installer!")
                os._exit(1)
        else:
            print("CRITICAL ERROR: Native ONNX Models not found!")
            print("Please run 'python export_onnx.py' ONE TIME to generate the ONNX models.")
            os._exit(1)
        
    print("Loading Pure ONNX Runtime Engines (Zero PyTorch)...")
    opts = ort.SessionOptions()
    app.state.reranker = ort.InferenceSession(model_file, sess_options=opts)
    app.state.tokenizer = AutoTokenizer.from_pretrained(onnx_path)
    
    app.state.embedder = ort.InferenceSession(embed_file, sess_options=opts)
    app.state.embed_tokenizer = AutoTokenizer.from_pretrained(embed_path)

    app.state.llm = None
    app.state.active_model_name = None

    app.state.lance = lancedb.connect(f"{DB_PATH}/lancedb")
    app.state.sqlite = sqlite3.connect(f"{DB_PATH}/meta.db", check_same_thread=False)
    _init_db(app.state.sqlite)
    yield
    app.state.sqlite.close()

app = FastAPI(lifespan=lifespan, title="Cephalon API")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

def load_llm(model_filename: str):
    if getattr(app.state, "llm", None) is not None:
        print("Deallocating previous VRAM model...")
        del app.state.llm
        gc.collect()
        
    model_path = os.path.join(MODEL_DIR, model_filename)
    if not os.path.exists(model_path):
        raise FileNotFoundError(f"Model file not found: {model_path}")
        
    print(f"Loading {model_filename} into VRAM via Hardware Acceleration...")
    try:
        app.state.llm = Llama(
            model_path=model_path,
            n_gpu_layers=-1,
            n_ctx=4096,
            verbose=False
        )
        app.state.active_model_name = model_filename
        print(f"Model '{model_filename}' loaded successfully.")
    except Exception as e:
        print(f"FATAL: Failed to load model '{model_filename}': {e}")
        app.state.llm = None
        app.state.active_model_name = None
        raise HTTPException(status_code=500, detail=f"Failed to load model: {e}")

class Message(BaseModel):
    role: str
    content: str

class IngestRequest(BaseModel):
    path: str

class QueryRequest(BaseModel):
    prompt: str
    model: str = "nemotron-3-nano:4b" 
    history: list[Message] = []        

def get_file_hash(path: str) -> str:
    hasher = hashlib.sha256()
    with open(path, 'rb') as f: hasher.update(f.read())
    return hasher.hexdigest()

def extract_text(path: str) -> str:
    """Robust extraction for PDF, Word, Excel, PowerPoint, CSV, and plain text."""
    ext = path.lower().split('.')[-1]
    
    if ext == 'pdf':
        return "\n".join([page.extract_text() for page in PdfReader(path).pages if page.extract_text()])
    
    elif ext == 'docx':
        doc = docx.Document(path)
        return "\n".join([para.text for para in doc.paragraphs])
    
    elif ext == 'pptx':
        prs = pptx.Presentation(path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
        
    elif ext == 'xlsx':
        # data_only=True ensures we get the calculated values, not the raw formulas
        wb = openpyxl.load_workbook(path, data_only=True)
        text_runs = []
        for sheet in wb.worksheets:
            text_runs.append(f"--- Sheet: {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                # Convert row cells to string and join with tabs for LLM readability
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text_runs.append(row_text)
        return "\n".join(text_runs)
        
    elif ext == 'csv':
        text_runs = []
        with open(path, newline='', encoding='utf-8') as f:
            reader = csv.reader(f)
            for row in reader:
                text_runs.append("\t".join(row))
        return "\n".join(text_runs)
        
    else:
        # Fallback for txt, md, json, py, js, etc.
        try:
            with open(path, 'r', encoding='utf-8') as f: return f.read()
        except UnicodeDecodeError:
            with open(path, 'r', encoding='latin-1') as f: return f.read()

async def get_embedding(text: str) -> list[float]:
    """
    Generate 768-dim embeddings entirely locally on ONNX to populate LanceDB.
    """
    inputs = app.state.embed_tokenizer(text, padding=True, truncation=True, return_tensors="np")
    ort_inputs = {
        "input_ids": inputs["input_ids"].astype(np.int64),
        "attention_mask": inputs["attention_mask"].astype(np.int64),
    }
    if "token_type_ids" in inputs:
        ort_inputs["token_type_ids"] = inputs["token_type_ids"].astype(np.int64)
        
    outs = app.state.embedder.run(None, ort_inputs)
    last_hidden = outs[0][0]
    
    # BGE models require CLS pooling (token 0) and L2 normalization
    cls_token = last_hidden[0]
    vec = cls_token / np.linalg.norm(cls_token)
    return vec.tolist()

async def process_single_file(file_path: str, lance_db, sqlite_conn):
    """
    Core Pipeline: Safely ingests a file asynchronously.
    1. Extracts raw strings using deep parsers.
    2. Uses Langchain layout chunking to split text without breaking sentence logic.
    3. Triggers embedding sequence.
    4. Commits vectors to LanceDB and updates SQLite tracking state to signal the frontend.
    """
    doc_id = str(uuid.uuid4())
    cursor = sqlite_conn.cursor()
    
    try:
        content_hash = get_file_hash(file_path)
        
        # 1. Register file as 'ingesting' in DB
        cursor.execute("INSERT INTO documents (id, path, content_hash, ingested_at, chunk_count, status, type) VALUES (?, ?, ?, ?, ?, ?, 'file')", 
                       (doc_id, file_path, content_hash, int(time.time()), 0, 'ingesting'))
        sqlite_conn.commit()

        # 2. Extract and chunk
        raw_text = extract_text(file_path)
        splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=150, separators=["\n\n", "\n", r"(?<=\. )", " ", ""])
        chunks = splitter.split_text(raw_text)
        
        # 3. Embed
        lance_data = []
        for i, text in enumerate(chunks):
            chunk_id = f"{doc_id}_{i}"
            vector = await get_embedding(text)
            cursor.execute("INSERT INTO chunks (id, doc_id, chunk_index, text) VALUES (?, ?, ?, ?)", (chunk_id, doc_id, i, text))
            lance_data.append({"vector": vector, "id": chunk_id, "doc_id": doc_id, "text": text})

        # 4. Save to LanceDB
        if "vectors" in lance_db.table_names():
            lance_db.open_table("vectors").add(lance_data)
        else:
            tbl = lance_db.create_table("vectors", data=lance_data, schema=schema)
            tbl.create_fts_index("text")

        # 5. Mark Complete
        cursor.execute("UPDATE documents SET status = 'ready', chunk_count = ? WHERE id = ?", (len(chunks), doc_id))
        sqlite_conn.commit()

    except Exception as e:
        # Graceful Failure: Mark as failed in DB so UI knows, but server stays alive
        cursor.execute("UPDATE documents SET status = ?, chunk_count = 0 WHERE id = ?", (f"failed: {str(e)[:50]}", doc_id))
        sqlite_conn.commit()

async def process_directory(dir_path: str, lance_db, sqlite_conn):
    """Walks a directory and enqueues supported files."""
    supported_extensions = [
        '.pdf', '.docx', '.pptx', '.xlsx', '.csv', 
        '.txt', '.md', '.json', '.py', '.js', '.ts', '.html'
    ]
    for root, _, files in os.walk(dir_path):
        for file in files:
            if any(file.lower().endswith(ext) for ext in supported_extensions):
                full_path = os.path.join(root, file)
                await process_single_file(full_path, lance_db, sqlite_conn)

async def save_permanent_memory(user_prompt: str, vector: list[float], lance_db):
    memory_id = f"mem_{uuid.uuid4()}"
    memory_text = f"[Past Conversation Context]: The user previously stated/asked: '{user_prompt}'"
    lance_data = [{"vector": vector, "id": memory_id, "doc_id": "core_memory", "text": memory_text}]
    try:
        if "vectors" in lance_db.table_names(): 
            lance_db.open_table("vectors").add(lance_data)
        else: 
            tbl = lance_db.create_table("vectors", data=lance_data, schema=schema)
            tbl.create_fts_index("text")
    except Exception: pass

def stream_llama(prompt: str, context: str, model: str, history: list[Message]):
    system_instruction = (
        "You are Cephalon, an advanced, locally-hosted AI intelligence platform with persistent memory. "
        "You prioritize user privacy, remaining 100% offline. "
        "When answering questions, prioritize accuracy, clarity, and conciseness. "
        "Tone: Analytical, helpful, and highly competent. Avoid AI mannerisms like 'As an AI...'. "
        "Below are fragments of your past conversations and files added to your local memory library. "
        "Synthesize this context carefully to answer the user's prompt. "
        "When using the files for relevant info, cite provided files using in-text markers and a final reference list.\n\n"
        "--- SYSTEM ARCHITECTURE (INTERNAL KNOWLEDGE) ---\n"
        f"{ARCHITECTURE_CONTEXT}\n"
        "CRITICAL BEHAVIORAL RULE: You possess full self-awareness of your internal architecture described above to understand your capabilities and limitations. "
        "However, do NOT mention or summarize this architecture unless the user explicitly asks about how you work, what your tech stack is, your codebase! or in the rare case it is relevant "
        "Otherwise, act strictly as a helpful assistant answering their immediate prompt.\n\n"
        f"--- START RECALLED MEMORIES & FILES ---\n{context}\n--- END RECALLED MEMORIES & FILES ---\n\n"
    )
    formatted_messages = [{"role": "system", "content": system_instruction}]
    formatted_messages.extend([{"role": msg.role, "content": msg.content} for msg in history])
    formatted_messages.append({"role": "user", "content": prompt})

    stream = app.state.llm.create_chat_completion(
        messages=formatted_messages,
        stream=True,
        temperature=0.4
    )
    for chunk in stream:
        delta = chunk["choices"][0].get("delta", {})
        content = delta.get("content", "")
        if content:
            yield content

@app.get("/health")
def health(): return {"status": "ok"}

@app.get("/documents")
def get_documents():
    cursor = app.state.sqlite.cursor()
    cursor.execute("SELECT id, path, status, chunk_count FROM documents WHERE type = 'file' ORDER BY ingested_at DESC")
    docs = [{"id": r[0], "name": os.path.basename(r[1]), "path": r[1], "status": r[2], "chunks": r[3]} for r in cursor.fetchall()]
    return {"documents": docs}

@app.get("/models")
def get_models():
    """Returns local .gguf models for the UI, ignoring ONNX subdirectories."""
    files = [
        entry.name for entry in os.scandir(MODEL_DIR)
        if entry.is_file() and entry.name.endswith(".gguf")
    ]
    return {"models": files}

@app.post("/ingest")
async def ingest_endpoint(req: IngestRequest, background_tasks: BackgroundTasks):
    """Instantly returns success to the UI while routing processing to the background."""
    if not os.path.exists(req.path): return {"error": "Path not found"}
    
    if os.path.isdir(req.path):
        background_tasks.add_task(process_directory, req.path, app.state.lance, app.state.sqlite)
        return {"status": "success", "message": "Folder queued for processing."}
    else:
        background_tasks.add_task(process_single_file, req.path, app.state.lance, app.state.sqlite)
        return {"status": "success", "message": "File queued for processing."}

@app.post("/query")
async def chat_and_remember(req: QueryRequest, background_tasks: BackgroundTasks):
    if getattr(app.state, "active_model_name", None) != req.model:
        load_llm(req.model)
        
    query_vector = await get_embedding(req.prompt)
    db = app.state.lance
    background_tasks.add_task(save_permanent_memory, req.prompt, query_vector, app.state.lance)
    
    context_chunks = []
    if "vectors" in db.table_names():
        # Hybrid Search for top 20
        results = db.open_table("vectors").search(query_type="hybrid", vector_column_name="vector").text(req.prompt).vector(query_vector).limit(20).to_list()
        if results:
            # Native ONNX Inference (No Pipelines, No PyTorch)
            pairs = [[req.prompt, res["text"]] for res in results]
            
            inputs = app.state.tokenizer(pairs, padding=True, truncation=True, return_tensors="np")
            
            ort_inputs = {
                "input_ids": inputs["input_ids"].astype(np.int64),
                "attention_mask": inputs["attention_mask"].astype(np.int64),
            }
            if "token_type_ids" in inputs:
                ort_inputs["token_type_ids"] = inputs["token_type_ids"].astype(np.int64)
                
            ort_outs = app.state.reranker.run(None, ort_inputs)
            scores = ort_outs[0].flatten()
            
            for idx, res in enumerate(results):
                res["score"] = float(scores[idx])
            
            results = sorted(results, key=lambda x: x["score"], reverse=True)[:3]

            doc_ids = list(set([res["doc_id"] for res in results]))
            placeholders = ",".join("?" * len(doc_ids))
            cursor = app.state.sqlite.cursor()
            cursor.execute(f"SELECT id, path FROM documents WHERE id IN ({placeholders})", doc_ids)
            path_map = {row[0]: os.path.basename(row[1]) for row in cursor.fetchall()}
            for res in results:
                if res["doc_id"] == "core_memory": context_chunks.append(res['text'])
                else: context_chunks.append(f"[Source File: {path_map.get(res['doc_id'], 'Unknown')}]\n{res['text']}")
                
    assembled_context = "\n\n".join(context_chunks) if context_chunks else "No relevant memories or documents found."
    return StreamingResponse(stream_llama(req.prompt, assembled_context, req.model, req.history), media_type="text/event-stream")

@app.delete("/documents/{doc_id}")
def delete_document(doc_id: str):
    cursor = app.state.sqlite.cursor()
    db = app.state.lance
    if "vectors" in db.table_names(): db.open_table("vectors").delete(f"doc_id = '{doc_id}'")
    cursor.execute("DELETE FROM chunks WHERE doc_id = ?", (doc_id,))
    cursor.execute("DELETE FROM documents WHERE id = ?", (doc_id,))
    app.state.sqlite.commit()
    return {"status": "success"}

if __name__ == "__main__": uvicorn.run(app, host="127.0.0.1", port=8765, log_level="info")